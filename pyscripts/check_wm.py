"""验证 wm_*.pptx 是否真的在每张幻灯片上注入了水印。"""
from pptx import Presentation
import sys

p = Presentation(sys.argv[1])
print("Slides:", len(p.slides))
wm_count = 0
wm_slides = []
for i, s in enumerate(p.slides):
    for sh in s.shapes:
        if sh.has_text_frame and "pptx-rs WATERMARK" in sh.text_frame.text:
            wm_count += 1
            wm_slides.append(i+1)
            break
print(f"Slides with watermark: {wm_count}/{len(p.slides)}")
print(f"First 5: {wm_slides[:5]}")

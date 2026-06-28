"""用 python-pptx 库（而非手动改 XML）添加水印，测试 WPS 兼容性。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

p = Presentation('_test/文旅IP人设打造抖音短视频运营方案.pptx')
print(f'Slides: {len(p.slides)}')

# 在每张幻灯片上添加一个文本框作为水印
for i, slide in enumerate(p.slides):
    # 添加文本框：左 1 英寸，上 4 英寸，宽 8 英寸，高 1 英寸
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(1)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p_obj = tf.paragraphs[0]
    p_obj.alignment = PP_ALIGN.CENTER
    run = p_obj.add_run()
    run.text = 'pptx-rs WATERMARK'
    # 红色 60pt 粗体
    font = run.font
    font.size = Pt(60)
    font.bold = True
    font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

p.save('_test_out/pptxlib_watermark.pptx')
print('pptxlib_watermark.pptx created')

# 对比 python-pptx 生成的 XML 和手动注入的 XML
import zipfile
with zipfile.ZipFile('_test_out/pptxlib_watermark.pptx') as z:
    slide1 = z.read('ppt/slides/slide1.xml').decode('utf-8')
    # 找到水印 shape
    wm_pos = slide1.find('pptx-rs WATERMARK')
    if wm_pos > 0:
        # 找到包含水印的 <p:sp> 元素
        # 向前找 <p:sp
        sp_start = slide1.rfind('<p:sp', 0, wm_pos)
        # 向后找 </p:sp>
        sp_end = slide1.find('</p:sp>', wm_pos) + len('</p:sp>')
        print()
        print('=== python-pptx generated watermark XML ===')
        print(slide1[sp_start:sp_end])

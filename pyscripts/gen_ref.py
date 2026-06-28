"""用 python-pptx 创建一个最小 hello.pptx，作为参考标准。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # 5 = Title Only

tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tb.text_frame.text = "Hello, python-pptx"
for p in tb.text_frame.paragraphs:
    for r in p.runs:
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x1F, 0x6F, 0xEB)

# 添加椭圆
from pptx.enum.shapes import MSO_SHAPE
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(3), Inches(3), Inches(2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
shape.line.fill.background()

# 第二张
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
tb2 = slide2.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(1))
tb2.text_frame.text = "第二张幻灯片"

prs.save("ref_hello.pptx")
print("已生成 ref_hello.pptx")
